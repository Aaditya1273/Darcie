"""
Darcie Presenton API — Self-contained styled presentation generator
Port: 8005

Does NOT require Presenton FastAPI to be running.
Generates real PPTX files using:
  1. Groq Llama 3.3 70B → structured slide content
  2. python-pptx → professional PPTX with real shapes, colors, fonts

Endpoints:
  POST /generate-report  { topic, context?, n_slides?, template? }
  GET  /download/{filename}
  GET  /health
"""

import os
import sys
import json
import asyncio
import uuid
from pathlib import Path
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
import httpx

load_dotenv(Path(__file__).parent / ".env")

GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "")
OUTPUT_DIR = Path(__file__).parent / "ppt_outputs"
OUTPUT_DIR.mkdir(exist_ok=True)

app = FastAPI(title="Darcie Presentation API", version="2.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])


class ReportRequest(BaseModel):
    topic: str
    context: str = ""
    n_slides: int = 8
    template: str = "professional"
    language: str = "English"
    tone: str = "professional"
    web_search: bool = False


class ReportResponse(BaseModel):
    presentation_id: str
    title: str
    view_url: str
    download_url: str
    slide_count: int


# ── Color themes ──────────────────────────────────────────────────
THEMES = {
    "professional": {
        "bg": (15, 23, 42),          # dark navy
        "accent": (59, 130, 246),    # blue
        "title_color": (255, 255, 255),
        "body_color": (203, 213, 225),
        "subtitle_color": (148, 163, 184),
    },
    "dark": {
        "bg": (17, 17, 17),
        "accent": (139, 92, 246),    # purple
        "title_color": (245, 245, 245),
        "body_color": (200, 200, 200),
        "subtitle_color": (150, 150, 150),
    },
    "green": {
        "bg": (5, 46, 22),
        "accent": (34, 197, 94),
        "title_color": (255, 255, 255),
        "body_color": (187, 247, 208),
        "subtitle_color": (134, 239, 172),
    },
}


async def _generate_slides_with_groq(topic: str, n_slides: int, context: str, tone: str) -> list[dict]:
    """Use Groq to generate structured slide content as JSON."""
    if not GROQ_API_KEY:
        raise RuntimeError("GROQ_API_KEY not set")

    system = """You are a professional presentation designer. Generate slide content as valid JSON.
Return ONLY a JSON array of slide objects, no other text.
Each slide object must have: "title" (string), "type" ("title"|"content"|"bullets"), "content" (string), "bullets" (array of strings, max 5).
For type "title": use for slide 1 only, set content as subtitle.
For type "bullets": use bullets array.
For type "content": use content string."""

    user = f"""Create a {n_slides}-slide {tone} presentation about: {topic}
{f'Context: {context[:2000]}' if context else ''}

Rules:
- Slide 1: type "title" with compelling title and subtitle
- Slides 2-{n_slides-1}: mix of "bullets" and "content" types
- Slide {n_slides}: type "bullets" with key takeaways
- All content must be specific and substantive, never placeholder text
- Return ONLY the JSON array"""

    async with httpx.AsyncClient(timeout=60.0) as client:
        r = await client.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
            json={
                "model": "llama-3.3-70b-versatile",
                "messages": [
                    {"role": "system", "content": system},
                    {"role": "user", "content": user},
                ],
                "temperature": 0.6,
                "max_tokens": 4000,
            }
        )
        if r.status_code != 200:
            raise RuntimeError(f"Groq error {r.status_code}: {r.text[:200]}")

        raw = r.json()["choices"][0]["message"]["content"].strip()

        # Strip markdown fences if present
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        raw = raw.strip()

        slides = json.loads(raw)
        if not isinstance(slides, list):
            raise RuntimeError("Groq did not return a JSON array")
        return slides


def _build_pptx(slides: list[dict], topic: str, theme_name: str = "professional") -> Path:
    """Build a real PPTX file using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    theme = THEMES.get(theme_name, THEMES["professional"])
    bg_color = RGBColor(*theme["bg"])
    accent_color = RGBColor(*theme["accent"])
    title_color = RGBColor(*theme["title_color"])
    body_color = RGBColor(*theme["body_color"])
    subtitle_color = RGBColor(*theme["subtitle_color"])

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Use blank layout for full control
    blank_layout = prs.slide_layouts[6]

    def add_bg(slide):
        """Fill slide background with theme color."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_text_box(slide, text, left, top, width, height, font_size, color, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        return txBox

    def add_accent_bar(slide, top_inches=0.08):
        """Add a thin accent color bar at top."""
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(top_inches),
            prs.slide_width, Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent_color
        bar.line.fill.background()

    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide)
        add_accent_bar(slide)

        slide_type = slide_data.get("type", "content")
        title_text = slide_data.get("title", f"Slide {i+1}")
        content_text = slide_data.get("content", "")
        bullets = slide_data.get("bullets", [])

        if slide_type == "title":
            # ── Title slide ──────────────────────────────────────
            # Large centered title
            add_text_box(slide, title_text,
                left=1.0, top=2.2, width=11.33, height=1.8,
                font_size=44, color=title_color, bold=True,
                align=PP_ALIGN.CENTER)
            # Subtitle
            if content_text:
                add_text_box(slide, content_text,
                    left=1.5, top=4.2, width=10.33, height=0.8,
                    font_size=20, color=subtitle_color,
                    align=PP_ALIGN.CENTER)
            # Topic label at bottom
            add_text_box(slide, topic.upper(),
                left=0.5, top=6.8, width=12.33, height=0.4,
                font_size=10, color=RGBColor(*theme["subtitle_color"]),
                align=PP_ALIGN.CENTER)

        elif slide_type == "bullets":
            # ── Bullet slide ─────────────────────────────────────
            # Title
            add_text_box(slide, title_text,
                left=0.6, top=0.4, width=11.5, height=0.9,
                font_size=28, color=title_color, bold=True)
            # Accent line under title
            line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.35), Inches(3.0), Inches(0.04))
            line.fill.solid(); line.fill.fore_color.rgb = accent_color; line.line.fill.background()

            # Bullets
            if bullets:
                txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(11.5), Inches(5.2))
                tf = txBox.text_frame
                tf.word_wrap = True
                for j, bullet in enumerate(bullets[:5]):
                    p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
                    p.space_before = Pt(8)
                    run = p.add_run()
                    run.text = f"  •  {bullet}"
                    run.font.size = Pt(18)
                    run.font.color.rgb = body_color
            elif content_text:
                add_text_box(slide, content_text,
                    left=0.6, top=1.6, width=11.5, height=5.2,
                    font_size=18, color=body_color)

        else:
            # ── Content slide ─────────────────────────────────────
            # Title
            add_text_box(slide, title_text,
                left=0.6, top=0.4, width=11.5, height=0.9,
                font_size=28, color=title_color, bold=True)
            # Accent line
            line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.35), Inches(3.0), Inches(0.04))
            line.fill.solid(); line.fill.fore_color.rgb = accent_color; line.line.fill.background()
            # Content
            add_text_box(slide, content_text or "\n".join(bullets),
                left=0.6, top=1.6, width=11.5, height=5.2,
                font_size=18, color=body_color)

        # Slide number
        add_text_box(slide, f"{i+1} / {len(slides)}",
            left=11.8, top=7.1, width=1.2, height=0.3,
            font_size=9, color=RGBColor(*theme["subtitle_color"]),
            align=PP_ALIGN.RIGHT)

    # Save
    filename = f"{uuid.uuid4().hex[:12]}.pptx"
    out_path = OUTPUT_DIR / filename
    prs.save(str(out_path))
    return out_path


# ── Routes ────────────────────────────────────────────────────────
@app.get("/health")
async def health():
    try:
        from pptx import Presentation
        pptx_ok = True
    except ImportError:
        pptx_ok = False
    return {
        "status": "ok",
        "engine": "groq+python-pptx (self-contained)",
        "groq_configured": bool(GROQ_API_KEY),
        "pptx_available": pptx_ok,
        "output_dir": str(OUTPUT_DIR),
    }


@app.post("/generate-report", response_model=ReportResponse)
@app.post("/generate-ppt", response_model=ReportResponse)
async def generate_report(req: ReportRequest):
    if not GROQ_API_KEY:
        raise HTTPException(503, "GROQ_API_KEY not configured")

    try:
        # 1. Generate slide content with Groq
        slides = await _generate_slides_with_groq(
            topic=req.topic,
            n_slides=req.n_slides,
            context=req.context,
            tone=req.tone,
        )
        print(f"[Presenton API] Generated {len(slides)} slides for: {req.topic}")

        # 2. Build PPTX
        pptx_path = _build_pptx(slides, req.topic, req.template)
        print(f"[Presenton API] PPTX saved: {pptx_path}")

        presentation_id = pptx_path.stem
        title = slides[0].get("title", req.topic) if slides else req.topic

        return ReportResponse(
            presentation_id=presentation_id,
            title=title,
            view_url=f"/download/{pptx_path.name}",
            download_url=f"/download/{pptx_path.name}",
            slide_count=len(slides),
        )

    except json.JSONDecodeError as e:
        raise HTTPException(500, f"Failed to parse slide content from LLM: {e}")
    except RuntimeError as e:
        raise HTTPException(500, str(e))
    except Exception as e:
        raise HTTPException(500, f"Presentation generation failed: {e}")


@app.get("/download/{filename}")
async def download(filename: str):
    # Security: only allow .pptx files from output dir
    if not filename.endswith(".pptx") or "/" in filename or ".." in filename:
        raise HTTPException(400, "Invalid filename")
    path = OUTPUT_DIR / filename
    if not path.exists():
        raise HTTPException(404, "File not found")
    return FileResponse(
        path=str(path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
    )


@app.get("/presentations")
async def list_presentations():
    files = list(OUTPUT_DIR.glob("*.pptx"))
    return [{"filename": f.name, "size": f.stat().st_size} for f in files]


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8005, log_level="info")
